from flask import Flask, request, jsonify, render_template, redirect, url_for
import os
import pandas as pd
from PIL import Image
import pytesseract
import PyPDF2
import json
import numpy as np
import logging
import re
import time

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Set Tesseract executable path
pytesseract.pytesseract.tesseract_cmd = '/usr/bin/tesseract'

# Function to process different file types
def extract_data(file_path, file_type):
    try:
        if file_type == 'csv':
            with open(file_path, 'r') as file:
                lines = file.readlines()
            addresses = [line.strip() for line in lines[1:] if line.strip()]
            return addresses

        elif file_type == 'xlsx':
            df = pd.read_excel(file_path)
            df = df.replace({np.nan: None})
            return df.to_dict(orient='records')

        elif file_type == 'pdf':
            pdf_data = []
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    pdf_data.append(page.extract_text())
            return " ".join(pdf_data)

        elif file_type in ['jpeg', 'jpg', 'png']:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text

        elif file_type in ['ppt', 'pptx']:
            import pptx
            ppt_data = []
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        ppt_data.append(shape.text)
            return " ".join([re.sub(r'\s+', ' ', slide.strip()) for slide in ppt_data if slide.strip()])

        else:
            return {'error': 'Unsupported file type'}

    except Exception as e:
        logger.error(f"Error extracting data from file: {e}")
        return {'error': f"Failed to process file: {e}"}

# Route for uploading and processing the file
@app.route('/')
def index():
    return '''
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Upload File</title>
    </head>
    <body>
        <h1>Upload a File</h1>
        <form action="/upload" method="post" enctype="multipart/form-data">
            <input type="file" name="file" required />
            <input type="submit" value="Upload" />
        </form>
    </body>
    </html>
    '''

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    # Save the file locally
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    file.save(file_path)

    # Extract data based on file type
    file_type = file.filename.rsplit('.', 1)[1].lower()
    extracted_data = extract_data(file_path, file_type)

    # Log extracted data
    logger.info(f"Extracted data: {extracted_data}")

    # Process data for table rendering (mock example)
    global combined_data_global
    combined_data_global = extracted_data  # Store for use in the table

    # Redirect to the table view
    return redirect(url_for('display_table'))

# Route to display extracted data in a table
@app.route('/table')
def display_table():
    global combined_data_global

    if not combined_data_global:
        return '''
        <h2>No data available</h2>
        <a href="/">Upload a new file</a>
        '''

    # Convert data to a pandas DataFrame and render as an HTML table
    if isinstance(combined_data_global, list):
        df = pd.DataFrame(combined_data_global)
    else:
        df = pd.DataFrame([combined_data_global])

    df = df.replace({np.nan: ''})  # Replace NaN values
    html_table = df.to_html(index=False, classes='table table-bordered table-hover', border=0)

    return f'''
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <title>Data Table</title>
        <link rel="stylesheet" href="https://cdn.datatables.net/1.11.5/css/jquery.dataTables.min.css">
        <script src="https://code.jquery.com/jquery-3.6.0.min.js"></script>
        <script src="https://cdn.datatables.net/1.11.5/js/jquery.dataTables.min.js"></script>
        <script>
            $(document).ready(function() {{
                $('.table').DataTable();
            }});
        </script>
    </head>
    <body>
        <h1>Extracted Data</h1>
        {html_table}
        <a href="/">Upload another file</a>
    </body>
    </html>
    '''

if __name__ == '__main__':
    app.run(debug=True)